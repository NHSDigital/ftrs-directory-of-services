#!/usr/bin/env python3
"""Generate a PPTX deck from the markdown slide file.

Reads docs/developer-guides/nfr-traceability-slides.md and converts sections delimited by '---' into slides.
Simple heuristics:
    - First non-empty line starting with '# ' => title (strip leading '#')
    - Else first non-empty line => title
    - Lines beginning with '-' or '*' become bullet points (stripped of leading marker)
    - Other non-empty lines become bullet points (to keep slides concise)

Output: docs/developer-guides/nfr-traceability-deck.pptx

Requires: python-pptx (install via `pip install python-pptx` or Make target).
"""
from pathlib import Path
import re
import sys

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed. Run: make deck-deps", file=sys.stderr)
    sys.exit(2)

# Defaults (can be overridden via CLI args: input_md [output_pptx])
DEFAULT_SLIDES_MD = Path('docs/developer-guides/nfr-traceability-slides.md')
DEFAULT_OUTPUT = Path('docs/developer-guides/nfr-traceability-deck.pptx')


def parse_sections(text: str):
    raw_sections = [s.strip() for s in re.split(r'\n---\n', text)]
    sections = []
    for sec in raw_sections:
        lines = [l.rstrip() for l in sec.splitlines() if l.strip()]
        if not lines:
            continue
        # Title detection
        title = None
        body_lines = []
        for i, line in enumerate(lines):
            if line.startswith('# '):
                title = line[2:].strip()
                body_lines = lines[i + 1 :]
                break
        if title is None:
            title = lines[0][:120].strip()
            body_lines = lines[1:]
        sections.append((title, body_lines))
    return sections


def to_bullets(body_lines):
    bullets = []
    for l in body_lines:
        # skip any redundant horizontal rules
        if l.strip() == '---':
            continue
        cleaned = re.sub(r'^[-*]\s*', '', l).strip()
        if cleaned:
            bullets.append(cleaned)
    return bullets[:15]  # safety cap per slide


def build_presentation(sections):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # First section becomes title slide
    if sections:
        title, body = sections[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if body:
            subtitle = '\n'.join(body[:3])
            slide.placeholders[1].text = subtitle[:300]

    for title, body in sections[1:]:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title[:120]
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        bullets = to_bullets(body)
        if not bullets:
            p = tf.paragraphs[0]
            p.text = '(No content)'
        else:
            for i, b in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = b[:300]
                p.level = 0
    return prs


def main():
    # CLI usage: script.py [input_markdown] [output_pptx]
    # If output not provided, derive name by replacing .md with .pptx or fall back to default.
    args = sys.argv[1:]
    if args and args[0] in {'-h', '--help'}:
        print("Usage: build_traceability_deck.py [input_markdown] [output_pptx]\n" \
            "Defaults: input='docs/developer-guides/nfr-traceability-slides.md' output='docs/developer-guides/nfr-traceability-deck.pptx'\n" \
            "Example condensed: build_traceability_deck.py docs/developer-guides/nfr-traceability-slides-condensed.md docs/developer-guides/nfr-traceability-deck-condensed.pptx")
        return 0

    if len(args) >= 1:
        slides_md = Path(args[0])
    else:
        slides_md = DEFAULT_SLIDES_MD

    if len(args) >= 2:
        output = Path(args[1])
    else:
        # derive output from input if input differs from default and endswith .md
        if slides_md != DEFAULT_SLIDES_MD and slides_md.suffix == '.md':
            output = slides_md.with_name(slides_md.stem + '.pptx')
        else:
            output = DEFAULT_OUTPUT

    if not slides_md.exists():
        print(f"ERROR: slides markdown not found: {slides_md}", file=sys.stderr)
        return 2

    text = slides_md.read_text(encoding='utf-8')
    sections = parse_sections(text)
    if not sections:
        print("ERROR: no sections parsed", file=sys.stderr)
        return 2
    prs = build_presentation(sections)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Deck written: {output}")
    return 0


if __name__ == '__main__':
    sys.exit(main())
